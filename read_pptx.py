from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

class ReadPptx:

    def __init__(self, path_pptx: str):
        self.path_pptx = path_pptx
    
    def extract_text_from_pptx_file(self):
        """
        Extract text from each slide in the PowerPoint file.

        This method reads the PowerPoint file specified by the `path_pptx` attribute,
        iterates through each slide, and extracts text from text boxes and placeholders.
        It stores the extracted text in the `slides_text` attribute.

        :return: A list of strings, where each string contains the text of a slide.
        """
        presentation = Presentation(self.path_pptx)
        self.slides_text = []

        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if shape.has_text_frame:
                        slide_text.append(shape.text_frame.text.strip())
            if slide_text:
                self.slides_text.append(' '.join(slide_text))
        return self.slides_text